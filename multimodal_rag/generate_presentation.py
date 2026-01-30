"""
PowerPoint Presentation Generator for Team ManageSphere
Table No. 18 - Multimodal RAG System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO
import os

class PresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Color scheme - Professional Blue & Orange
        self.primary_color = RGBColor(30, 58, 95)  # Dark Blue
        self.accent_color = RGBColor(240, 128, 48)  # Orange
        self.text_color = RGBColor(51, 51, 51)     # Dark Gray
        self.light_bg = RGBColor(248, 249, 250)    # Light Gray
        
    def add_title_slide(self):
        """Slide 1: Title Slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Background
        background = slide.shapes.add_shape(
            1, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.primary_color
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "🧠 Multimodal RAG System"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5), Inches(9), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Evidence-Based Multimodal Retrieval & Generation"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = self.accent_color
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Team Info
        team_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(9), Inches(1)
        )
        team_frame = team_box.text_frame
        team_frame.text = "Team: ManageSphere | Table No. 18"
        team_para = team_frame.paragraphs[0]
        team_para.font.size = Pt(20)
        team_para.font.color.rgb = RGBColor(255, 255, 255)
        team_para.alignment = PP_ALIGN.CENTER
        
    def add_problem_statement_slide(self):
        """Slide 2: Problem Statement"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header(slide, "Problem Statement & Solution")
        
        # Problem section
        problem_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5)
        )
        problem_frame = problem_box.text_frame
        problem_frame.word_wrap = True
        
        p = problem_frame.add_paragraph()
        p.text = "❌ Current Challenges:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        
        challenges = [
            "Information scattered across multiple formats",
            "Inability to query multimodal data uniformly",
            "Lack of evidence tracing and conflict detection",
            "Hallucinations in AI responses"
        ]
        
        for challenge in challenges:
            p = problem_frame.add_paragraph()
            p.text = f"• {challenge}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.text_color
            p.level = 1
            
        # Solution section
        solution_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(1.5), Inches(4), Inches(2.5)
        )
        solution_frame = solution_box.text_frame
        solution_frame.word_wrap = True
        
        s = solution_frame.add_paragraph()
        s.text = "✅ Our Solution:"
        s.font.size = Pt(18)
        s.font.bold = True
        s.font.color.rgb = self.accent_color
        
        solutions = [
            "Unified multimodal knowledge base",
            "Cross-modal semantic retrieval",
            "Evidence-grounded responses",
            "Built-in conflict detection",
            "Uncertainty-aware AI"
        ]
        
        for solution in solutions:
            s = solution_frame.add_paragraph()
            s.text = f"✓ {solution}"
            s.font.size = Pt(14)
            s.font.color.rgb = self.text_color
            s.level = 1
            
    def add_business_model_slide(self):
        """Slide 3: Business Model (MANDATORY)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header(slide, "Business Model")
        
        # Target Markets
        market_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.8)
        )
        market_frame = market_box.text_frame
        market_frame.word_wrap = True
        
        m = market_frame.add_paragraph()
        m.text = "🎯 Target Markets:"
        m.font.size = Pt(16)
        m.font.bold = True
        m.font.color.rgb = self.primary_color
        
        markets = [
            "Enterprise Knowledge Management",
            "Legal & Compliance Firms",
            "Healthcare Documentation",
            "Research Institutions",
            "Educational Platforms"
        ]
        
        for market in markets:
            m = market_frame.add_paragraph()
            m.text = f"• {market}"
            m.font.size = Pt(13)
            m.font.color.rgb = self.text_color
            m.level = 1
            
        # Revenue Streams
        revenue_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(1.5), Inches(4), Inches(2.8)
        )
        revenue_frame = revenue_box.text_frame
        revenue_frame.word_wrap = True
        
        r = revenue_frame.add_paragraph()
        r.text = "💰 Revenue Streams:"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = self.accent_color
        
        revenues = [
            "SaaS Subscription (Tiered)",
            "Enterprise Licensing",
            "API Usage Fees",
            "Custom Integration Services",
            "Premium Support Plans"
        ]
        
        for revenue in revenues:
            r = revenue_frame.add_paragraph()
            r.text = f"💵 {revenue}"
            r.font.size = Pt(13)
            r.font.color.rgb = self.text_color
            r.level = 1
            
        # Value Proposition
        value_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.7), Inches(9), Inches(2)
        )
        value_frame = value_box.text_frame
        value_frame.word_wrap = True
        
        v = value_frame.add_paragraph()
        v.text = "⭐ Unique Value Proposition:"
        v.font.size = Pt(16)
        v.font.bold = True
        v.font.color.rgb = self.primary_color
        
        v = value_frame.add_paragraph()
        v.text = "First evidence-grounded multimodal RAG system with built-in conflict detection, uncertainty awareness, and hallucination prevention. Self-hostable for data privacy."
        v.font.size = Pt(13)
        v.font.color.rgb = self.text_color
        
    def add_key_features_slide(self):
        """Slide 4: Key Features"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header(slide, "Key Features & Capabilities")
        
        features = [
            ("📄 Multimodal Ingestion", "PDF, DOCX, TXT, Images, Audio"),
            ("🔍 Cross-Modal Retrieval", "Unified semantic search across all formats"),
            ("🎯 Evidence Grounding", "All responses cite specific sources"),
            ("⚠️ Conflict Detection", "Identifies contradictory information"),
            ("🤔 Uncertainty Awareness", "Explicitly acknowledges knowledge gaps"),
            ("🚫 Hallucination Prevention", "Refuses to answer without evidence"),
            ("🌍 30+ Languages", "Auto-translation for global access"),
            ("☁️ Cloud Integration", "Google Drive, OneDrive, S3, Dropbox")
        ]
        
        # Create two columns
        left_x, right_x = Inches(0.5), Inches(5.25)
        y_start = Inches(1.5)
        box_height = Inches(0.65)
        
        for i, (feature, description) in enumerate(features):
            x = left_x if i < 4 else right_x
            y = y_start + (i % 4) * box_height
            
            # Feature box
            box = slide.shapes.add_textbox(x, y, Inches(4.5), box_height)
            frame = box.text_frame
            frame.word_wrap = True
            frame.margin_top = Inches(0.05)
            
            # Feature title
            p = frame.paragraphs[0]
            p.text = feature
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            
            # Description
            p = frame.add_paragraph()
            p.text = description
            p.font.size = Pt(11)
            p.font.color.rgb = self.text_color
            
    def add_architecture_slide(self):
        """Slide 5: System Architecture (MANDATORY)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header(slide, "System Architecture")
        
        # Architecture components
        components = [
            ("Frontend Layer", [
                "Web UI (HTML/CSS/JS)",
                "Voice Input Interface",
                "Real-time Sync Dashboard"
            ]),
            ("API Layer", [
                "FastAPI REST Endpoints",
                "/upload, /query, /summarize",
                "WebSocket for live updates"
            ]),
            ("Processing Layer", [
                "Text/Image/Audio Processors",
                "OCR (EasyOCR + Tesseract)",
                "Whisper ASR for audio"
            ]),
            ("Embedding Layer", [
                "Text: MiniLM-L6 (384-dim)",
                "Images: CLIP ViT-B/32",
                "Cross-modal alignment"
            ]),
            ("Storage Layer", [
                "ChromaDB Vector Store",
                "HNSW Index for speed",
                "Metadata filtering"
            ]),
            ("Generation Layer", [
                "LLM: Ollama/Gemini",
                "RAG Pipeline",
                "Confidence scoring"
            ])
        ]
        
        # Create grid layout
        boxes_per_row = 3
        box_width = Inches(3)
        box_height = Inches(1.8)
        x_start = Inches(0.5)
        y_start = Inches(1.5)
        x_gap = Inches(0.25)
        y_gap = Inches(0.3)
        
        for i, (title, items) in enumerate(components):
            row = i // boxes_per_row
            col = i % boxes_per_row
            
            x = x_start + col * (box_width + x_gap)
            y = y_start + row * (box_height + y_gap)
            
            # Add box with background
            box = slide.shapes.add_shape(
                1,  # Rectangle
                x, y, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.light_bg
            box.line.color.rgb = self.primary_color
            box.line.width = Pt(2)
            
            # Add text
            text_box = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(0.1),
                box_width - Inches(0.2), box_height - Inches(0.2)
            )
            frame = text_box.text_frame
            frame.word_wrap = True
            frame.margin_top = Inches(0.05)
            
            # Title
            p = frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            p.alignment = PP_ALIGN.CENTER
            
            # Items
            for item in items:
                p = frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(9)
                p.font.color.rgb = self.text_color
                p.level = 1
                
    def add_process_flow_slide(self):
        """Slide 6: Process Flow Chart (MANDATORY)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header(slide, "Process Flow Chart")
        
        # Define flow steps
        flows = [
            ("1. Document Upload", "User uploads PDF/DOCX/Image/Audio"),
            ("2. Content Extraction", "Text/OCR/Speech-to-Text processing"),
            ("3. Chunking", "Split into semantic chunks"),
            ("4. Embedding Generation", "Convert to 384-dim vectors"),
            ("5. Vector Storage", "Store in ChromaDB with metadata"),
            ("6. Query Processing", "User asks question (text/voice)"),
            ("7. Retrieval", "Semantic search across modalities"),
            ("8. Conflict Detection", "Check for contradictions"),
            ("9. Response Generation", "LLM generates evidence-based answer"),
            ("10. Confidence Scoring", "Calculate reliability score")
        ]
        
        # Create vertical flow
        box_width = Inches(8.5)
        box_height = Inches(0.45)
        x = Inches(0.75)
        y_start = Inches(1.5)
        y_gap = Inches(0.08)
        
        for i, (step, description) in enumerate(flows):
            y = y_start + i * (box_height + y_gap)
            
            # Alternating colors
            bg_color = self.primary_color if i % 2 == 0 else self.accent_color
            
            # Step box
            box = slide.shapes.add_shape(1, x, y, box_width, box_height)
            box.fill.solid()
            box.fill.fore_color.rgb = bg_color
            box.line.fill.background()
            
            # Text
            text_box = slide.shapes.add_textbox(
                x + Inches(0.1), y,
                box_width - Inches(0.2), box_height
            )
            frame = text_box.text_frame
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = frame.paragraphs[0]
            p.text = f"{step}: {description}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True if i in [0, 5, 9] else False
            
            # Add arrow between steps (except last)
            if i < len(flows) - 1:
                arrow_y = y + box_height + Inches(0.01)
                arrow = slide.shapes.add_shape(
                    87,  # Down arrow
                    x + box_width / 2 - Inches(0.15), arrow_y,
                    Inches(0.3), Inches(0.06)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.text_color
                arrow.line.fill.background()
                
    def add_tech_stack_slide(self):
        """Slide 7: Technology Stack"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header(slide, "Technology Stack")
        
        tech_categories = [
            ("🎨 Frontend", ["HTML5 / CSS3", "Vanilla JavaScript", "Progressive Web App"]),
            ("⚙️ Backend", ["FastAPI (Python)", "Async/Await", "RESTful APIs"]),
            ("🤖 AI/ML Models", ["Llama 3.2 (3B)", "CLIP ViT-B/32", "Whisper Tiny"]),
            ("📊 Embeddings", ["Sentence Transformers", "MiniLM-L6-v2", "384-dim vectors"]),
            ("💾 Storage", ["ChromaDB", "HNSW Index", "SQLite metadata"]),
            ("📄 Processing", ["PyPDF / python-docx", "EasyOCR / Tesseract", "Pydub audio"]),
            ("☁️ Integration", ["Google Drive API", "OneDrive API", "AWS S3 SDK"]),
            ("🌐 Translation", ["Langdetect", "Googletrans", "30+ languages"])
        ]
        
        # Create 4x2 grid
        box_width = Inches(4.5)
        box_height = Inches(1.4)
        x_start = Inches(0.5)
        y_start = Inches(1.5)
        x_gap = Inches(0.5)
        y_gap = Inches(0.25)
        
        for i, (category, techs) in enumerate(tech_categories):
            row = i // 2
            col = i % 2
            
            x = x_start + col * (box_width + x_gap)
            y = y_start + row * (box_height + y_gap)
            
            # Category box
            box = slide.shapes.add_textbox(x, y, box_width, box_height)
            frame = box.text_frame
            frame.word_wrap = True
            
            # Category title
            p = frame.paragraphs[0]
            p.text = category
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            
            # Technologies
            for tech in techs:
                p = frame.add_paragraph()
                p.text = f"• {tech}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.text_color
                p.level = 1
                
    def add_competitive_advantage_slide(self):
        """Slide 8: Competitive Advantage"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header(slide, "Competitive Advantage")
        
        # Our System vs Others
        comparison_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        )
        frame = comparison_box.text_frame
        frame.word_wrap = True
        
        # Title
        p = frame.paragraphs[0]
        p.text = "🏆 Why Choose Our System?"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.alignment = PP_ALIGN.CENTER
        
        advantages = [
            ("✅ Evidence Provenance", "Every answer cites exact sources - no hallucinations"),
            ("✅ Conflict Detection", "Automatically identifies contradictory information"),
            ("✅ Uncertainty Handling", "Transparently admits knowledge gaps"),
            ("✅ True Multimodal", "Seamless integration of text, images, and audio"),
            ("✅ Self-Hostable", "Full data control - no vendor lock-in"),
            ("✅ Offline Capable", "Works without internet connectivity"),
            ("✅ Multilingual Support", "30+ languages with auto-detection"),
            ("✅ Production Ready", "Complete system with PWA, cloud sync, real-time updates")
        ]
        
        for adv_title, adv_desc in advantages:
            p = frame.add_paragraph()
            p.text = adv_title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.accent_color
            p.space_after = Pt(2)
            
            p = frame.add_paragraph()
            p.text = adv_desc
            p.font.size = Pt(12)
            p.font.color.rgb = self.text_color
            p.level = 1
            p.space_after = Pt(8)
            
    def add_use_cases_slide(self):
        """Slide 9: Real-World Use Cases"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header(slide, "Real-World Use Cases")
        
        use_cases = [
            ("⚖️ Legal Research", [
                "Query across case files, precedents, contracts",
                "Detect conflicting clauses automatically",
                "Evidence-backed legal analysis"
            ]),
            ("🏥 Healthcare", [
                "Medical records + imaging + audio notes",
                "Patient history with full traceability",
                "Multi-source diagnosis support"
            ]),
            ("🏢 Enterprise Knowledge", [
                "Company docs, presentations, meeting recordings",
                "Cross-department information retrieval",
                "Policy compliance checking"
            ]),
            ("🎓 Academic Research", [
                "Papers, datasets, lecture recordings",
                "Literature review automation",
                "Contradiction detection in studies"
            ])
        ]
        
        # 2x2 grid
        box_width = Inches(4.5)
        box_height = Inches(2.5)
        x_start = Inches(0.5)
        y_start = Inches(1.5)
        x_gap = Inches(0.5)
        y_gap = Inches(0.4)
        
        for i, (title, points) in enumerate(use_cases):
            row = i // 2
            col = i % 2
            
            x = x_start + col * (box_width + x_gap)
            y = y_start + row * (box_height + y_gap)
            
            # Background box
            bg_box = slide.shapes.add_shape(1, x, y, box_width, box_height)
            bg_box.fill.solid()
            bg_box.fill.fore_color.rgb = self.light_bg
            bg_box.line.color.rgb = self.primary_color
            bg_box.line.width = Pt(1.5)
            
            # Text box
            text_box = slide.shapes.add_textbox(
                x + Inches(0.15), y + Inches(0.15),
                box_width - Inches(0.3), box_height - Inches(0.3)
            )
            frame = text_box.text_frame
            frame.word_wrap = True
            
            # Title
            p = frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            
            # Points
            for point in points:
                p = frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.text_color
                p.level = 1
                
    def add_thank_you_slide(self):
        """Slide 10: Thank You"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.primary_color
        background.line.fill.background()
        
        # Thank you message
        thank_you_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1)
        )
        frame = thank_you_box.text_frame
        p = frame.paragraphs[0]
        p.text = "Thank You!"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Team info
        team_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8), Inches(9), Inches(0.6)
        )
        frame = team_box.text_frame
        p = frame.paragraphs[0]
        p.text = "Team ManageSphere | Table No. 18"
        p.font.size = Pt(24)
        p.font.color.rgb = self.accent_color
        p.alignment = PP_ALIGN.CENTER
        
        # Contact
        contact_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(9), Inches(1)
        )
        frame = contact_box.text_frame
        
        p = frame.paragraphs[0]
        p.text = "GitHub: github.com/rahulranjan9770/multimodel"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = frame.add_paragraph()
        p.text = "Evidence-Based Multimodal RAG System"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
        
    def _add_header(self, slide, title):
        """Helper: Add consistent header to slides"""
        # Header background
        header_bg = slide.shapes.add_shape(
            1, 0, 0, self.prs.slide_width, Inches(1)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.primary_color
        header_bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(8), Inches(0.6)
        )
        frame = title_box.text_frame
        p = frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Team badge
        badge_box = slide.shapes.add_textbox(
            Inches(8.5), Inches(0.25), Inches(1.3), Inches(0.5)
        )
        badge_frame = badge_box.text_frame
        badge_para = badge_frame.paragraphs[0]
        badge_para.text = "Table 18"
        badge_para.font.size = Pt(12)
        badge_para.font.color.rgb = self.accent_color
        badge_para.alignment = PP_ALIGN.CENTER
        
    def generate(self, output_path="ManageSphere_Presentation.pptx"):
        """Generate the complete presentation"""
        print("🎨 Generating presentation slides...")
        
        # Add all slides in order
        self.add_title_slide()                    # Slide 1
        self.add_problem_statement_slide()        # Slide 2
        self.add_business_model_slide()           # Slide 3 - MANDATORY
        self.add_key_features_slide()             # Slide 4
        self.add_architecture_slide()             # Slide 5 - MANDATORY
        self.add_process_flow_slide()             # Slide 6 - MANDATORY
        self.add_tech_stack_slide()               # Slide 7
        self.add_competitive_advantage_slide()    # Slide 8
        self.add_use_cases_slide()                # Slide 9
        self.add_thank_you_slide()                # Slide 10
        
        # Save presentation
        self.prs.save(output_path)
        print(f"✅ Presentation saved: {output_path}")
        print(f"📊 Total slides: {len(self.prs.slides)}")
        
        return output_path


if __name__ == "__main__":
    generator = PresentationGenerator()
    output_file = generator.generate()
    print(f"\n🎉 Presentation ready: {output_file}")
    print("\n📋 Slide Breakdown:")
    print("  1. Title Slide")
    print("  2. Problem Statement & Solution")
    print("  3. Business Model ⭐ (MANDATORY)")
    print("  4. Key Features & Capabilities")
    print("  5. System Architecture ⭐ (MANDATORY)")
    print("  6. Process Flow Chart ⭐ (MANDATORY)")
    print("  7. Technology Stack")
    print("  8. Competitive Advantage")
    print("  9. Real-World Use Cases")
    print(" 10. Thank You")
