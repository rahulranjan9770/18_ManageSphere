"""Automated Presentation Generator using knowledge base context."""
import io
import os
import time
import uuid
import re
from datetime import datetime
from typing import List, Optional, Dict, Any, Tuple
from pathlib import Path
from pydantic import BaseModel, Field
from enum import Enum

# Direct imports - python-pptx is required
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from backend.retrieval.cross_modal_retriever import CrossModalRetriever
from backend.generation.llm_client import LLMClient
from backend.config import settings
from backend.utils.logger import logger


class PresentationTheme(str, Enum):
    """Available presentation themes."""
    PROFESSIONAL = "professional"
    MODERN = "modern"
    MINIMAL = "minimal"
    CORPORATE = "corporate"
    CREATIVE = "creative"


class SlideType(str, Enum):
    """Types of slides that can be generated."""
    TITLE = "title"
    BULLET_POINTS = "bullet_points"
    TWO_COLUMN = "two_column"
    CONTENT = "content"
    QUOTE = "quote"
    SUMMARY = "summary"


class SlideContent(BaseModel):
    """Content for a single slide."""
    slide_type: SlideType
    title: str
    subtitle: Optional[str] = None
    bullet_points: List[str] = Field(default_factory=list)
    content: Optional[str] = None
    left_column: List[str] = Field(default_factory=list)
    right_column: List[str] = Field(default_factory=list)
    source_references: List[str] = Field(default_factory=list)
    notes: Optional[str] = None


class PresentationRequest(BaseModel):
    """Request model for generating presentations."""
    topic: str  # Main topic or goal
    source_documents: List[str] = Field(default_factory=list)  # Specific docs to use
    num_slides: int = 5  # Target number of slides
    theme: PresentationTheme = PresentationTheme.PROFESSIONAL
    include_title_slide: bool = True
    include_summary_slide: bool = True
    include_sources: bool = True
    additional_instructions: Optional[str] = None


class PresentationResponse(BaseModel):
    """Response model for generated presentations."""
    presentation_id: str
    filename: str
    topic: str
    num_slides: int
    slides: List[SlideContent]
    sources_used: List[str]
    theme: str
    processing_time: float
    download_url: str


# Theme color configurations
THEME_COLORS = {
    "professional": {
        "primary": RGBColor(0x1E, 0x3A, 0x5F),  # Dark blue
        "secondary": RGBColor(0x3D, 0x5A, 0x80),
        "accent": RGBColor(0xF0, 0x80, 0x30),  # Orange
        "text": RGBColor(0x33, 0x33, 0x33),
        "background": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "modern": {
        "primary": RGBColor(0x6C, 0x5C, 0xE7),  # Purple
        "secondary": RGBColor(0xA2, 0x9B, 0xFE),
        "accent": RGBColor(0x00, 0xD2, 0xD3),  # Teal
        "text": RGBColor(0x2D, 0x3A, 0x4B),
        "background": RGBColor(0xF8, 0xF9, 0xFA),
    },
    "minimal": {
        "primary": RGBColor(0x2D, 0x3A, 0x4B),  # Dark gray
        "secondary": RGBColor(0x6C, 0x75, 0x7D),
        "accent": RGBColor(0x00, 0x7B, 0xFF),  # Blue
        "text": RGBColor(0x21, 0x25, 0x29),
        "background": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "corporate": {
        "primary": RGBColor(0x00, 0x52, 0x8A),  # Corporate blue
        "secondary": RGBColor(0x00, 0x79, 0xC1),
        "accent": RGBColor(0xFF, 0xB8, 0x00),  # Gold
        "text": RGBColor(0x33, 0x33, 0x33),
        "background": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "creative": {
        "primary": RGBColor(0xE0, 0x43, 0x89),  # Pink
        "secondary": RGBColor(0xFF, 0x6B, 0x6B),
        "accent": RGBColor(0x4E, 0xCE, 0xC3),  # Teal
        "text": RGBColor(0x2F, 0x2F, 0x2F),
        "background": RGBColor(0xFF, 0xFA, 0xFC),
    }
}


class PresentationGenerator:
    """Automated PowerPoint presentation generator using RAG context."""
    
    def __init__(self, retriever: CrossModalRetriever):
        self.retriever = retriever
        self.llm_client = LLMClient()
        self.output_dir = settings.upload_dir / "presentations"
        self.output_dir.mkdir(exist_ok=True)
    
    async def generate_presentation(self, request: PresentationRequest) -> PresentationResponse:
        """Generate a complete PowerPoint presentation."""
        start_time = time.time()
        presentation_id = str(uuid.uuid4())[:8]
        
        # Step 1: Retrieve relevant context from knowledge base
        sources = await self._retrieve_context(request.topic, request.source_documents)
        sources_used = list(set(s.source_file for s in sources))
        
        # Step 2: Generate slide content using LLM
        slides_content = await self._generate_slides_content(request, sources)
        
        # Step 3: Create the PowerPoint presentation
        pptx_path, filename = self._create_pptx(
            slides_content, 
            request.theme,
            presentation_id
        )
        
        processing_time = time.time() - start_time
        
        logger.logger.info(
            f"📊 Generated presentation: {filename} with {len(slides_content)} slides "
            f"in {processing_time:.1f}s"
        )
        
        return PresentationResponse(
            presentation_id=presentation_id,
            filename=filename,
            topic=request.topic,
            num_slides=len(slides_content),
            slides=slides_content,
            sources_used=sources_used,
            theme=request.theme.value,
            processing_time=processing_time,
            download_url=f"/presentations/{filename}"
        )
    
    async def _retrieve_context(
        self, 
        topic: str, 
        specific_docs: List[str]
    ) -> List:
        """Retrieve relevant context from knowledge base."""
        try:
            sources = self.retriever.retrieve(query=topic, top_k=8)  # Reduced from 15 to 8 for speed
            
            # Prioritize specific documents if mentioned
            if specific_docs:
                specific_matches = []
                other_matches = []
                
                for source in sources:
                    source_lower = source.source_file.lower()
                    if any(doc.lower() in source_lower for doc in specific_docs):
                        specific_matches.append(source)
                    else:
                        other_matches.append(source)
                
                sources = specific_matches + other_matches
            
            return sources
        except Exception as e:
            logger.logger.error(f"Context retrieval error: {e}")
            return []
    
    async def _generate_slides_content(
        self, 
        request: PresentationRequest, 
        sources: List
    ) -> List[SlideContent]:
        """Generate slide content using LLM with enhanced content extraction."""
        slides = []
        
        # Build richer evidence text with key points extraction
        evidence_parts = []
        for i, source in enumerate(sources[:5], 1):
            # Extract more meaningful chunks
            content = source.content[:400]  # Slightly more content
            evidence_parts.append(
                f"Document {i} ({source.source_file}):\n{content}"
            )
        evidence_text = "\n\n".join(evidence_parts) if evidence_parts else "No documents available."
        
        # Calculate content slides needed
        content_slides = request.num_slides
        if request.include_title_slide:
            content_slides -= 1
        if request.include_summary_slide:
            content_slides -= 1
        content_slides = max(1, content_slides)
        
        # ENHANCED Prompt with explicit content requirements
        prompt = f"""You are an expert presentation creator. Create {content_slides} detailed slides about "{request.topic}".

REFERENCE MATERIAL:
{evidence_text}

{f"SPECIAL INSTRUCTIONS: {request.additional_instructions}" if request.additional_instructions else ""}

Create EXACTLY {content_slides} slides. For each slide, use this format:

---SLIDE---
TITLE: [Specific, descriptive title]
BULLETS: 
• [First key point with specific details or data]
• [Second point with actionable insight]
• [Third point with supporting information]
• [Fourth point with concrete example or benefit]
• [Fifth point if relevant]
---END---

REQUIREMENTS:
- Each bullet point must be 10-15 words long
- Include specific facts, numbers, or examples from the documents
- Make points actionable and clear
- Use professional language
- Each slide should cover a distinct aspect of the topic

Generate all {content_slides} slides now with substantive content:"""

        try:
            response = await self.llm_client.generate(
                prompt=prompt,
                max_tokens=1000,  # Increased slightly for more content
                temperature=0.5  # Balanced creativity and consistency
            )
            
            # Parse the LLM response with enhanced parsing
            slides = self._parse_slides_response(response, request)
            
            # Enhance parsed slides with more content if needed
            slides = self._enhance_slide_content(slides, sources)
            
            # If still no slides, use better fallback
            if not slides or len(slides) < content_slides:
                logger.logger.warning(f"Only {len(slides)} slides parsed, using fallback")
                slides = self._create_fallback_slides(request, sources)
            
        except Exception as e:
            logger.logger.error(f"LLM generation error: {e}")
            slides = self._create_fallback_slides(request, sources)
        
        # Add title slide if requested
        if request.include_title_slide:
            title_slide = SlideContent(
                slide_type=SlideType.TITLE,
                title=request.topic,
                subtitle=f"Evidence-Based AI Presentation",
                notes=f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            )
            slides.insert(0, title_slide)
        
        # Add summary slide if requested
        if request.include_summary_slide:
            # Create a more detailed summary
            source_names = list(set(s.source_file for s in sources[:3]))
            summary_points = [
                f"✅ Key insights drawn from {len(sources)} knowledge base chunks",
                f"📚 Primary sources: {', '.join(source_names)}",
                f"📊 {len(slides)} detailed content slides created",
                "💡 For in-depth details, consult original documents"
            ]
            summary_slide = SlideContent(
                slide_type=SlideType.SUMMARY,
                title="Key Takeaways & Next Steps",
                bullet_points=summary_points,
                notes="Summary of main findings"
            )
            slides.append(summary_slide)
        
        return slides
    
    def _enhance_slide_content(self, slides: List[SlideContent], sources: List) -> List[SlideContent]:
        """Enhance slides that have minimal content."""
        enhanced_slides = []
        
        for slide in slides:
            # If slide has fewer than 3 bullets, try to add more
            if len(slide.bullet_points) < 3 and sources:
                # Extract additional relevant points from sources
                for source in sources[:2]:
                    if len(slide.bullet_points) >= 5:
                        break
                    
                    # Extract sentences from source
                    sentences = [s.strip() for s in source.content.split('.') if len(s.strip()) > 30]
                    for sent in sentences[:2]:
                        if len(slide.bullet_points) < 5 and len(sent) < 100:
                            slide.bullet_points.append(sent + ('.' if not sent.endswith('.') else ''))
            
            enhanced_slides.append(slide)
        
        return enhanced_slides
    
    def _parse_slides_response(
        self, 
        response: str, 
        request: PresentationRequest
    ) -> List[SlideContent]:
        """Parse LLM response into slide content objects with improved bullet extraction."""
        slides = []
        
        # Split by slide markers
        slide_blocks = re.split(r'---SLIDE---', response)
        
        for block in slide_blocks:
            if not block.strip() or '---END---' not in block:
                continue
            
            block = block.split('---END---')[0].strip()
            
            # Extract components
            slide_type = SlideType.BULLET_POINTS
            title = "Untitled Slide"
            bullets = []
            left_col = []
            right_col = []
            sources = []
            
            lines = block.split('\n')
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                
                if line.startswith('TYPE:'):
                    type_str = line[5:].strip().lower()
                    if 'two_column' in type_str:
                        slide_type = SlideType.TWO_COLUMN
                    elif 'content' in type_str:
                        slide_type = SlideType.CONTENT
                    else:
                        slide_type = SlideType.BULLET_POINTS
                        
                elif line.startswith('TITLE:'):
                    title = line[6:].strip()
                    
                elif line.startswith('BULLETS:'):
                    # Handle different bullet formats
                    rest = line[8:].strip()
                    
                    # Format 1: Pipe-separated on same line
                    if '|' in rest:
                        bullets = [b.strip() for b in rest.split('|') if b.strip()]
                    # Format 2: Multi-line bullets following
                    elif not rest:
                        i += 1
                        while i < len(lines):
                            next_line = lines[i].strip()
                            if next_line.startswith(('•', '-', '*', '–')) or (next_line and next_line[0].isdigit() and next_line[1:3] in ['. ', ') ']):
                                # Remove bullet markers
                                bullet_text = re.sub(r'^[•\-\*–]\s*|^\d+[.)]\s*', '', next_line).strip()
                                if bullet_text and len(bullet_text) > 10:  # Meaningful content
                                    bullets.append(bullet_text)
                            elif next_line and not next_line.startswith(('TITLE', 'TYPE', 'LEFT', 'RIGHT', 'SOURCES', '---')):
                                # Treat as continuation or new bullet
                                if next_line:
                                    bullets.append(next_line)
                            else:
                                i -= 1  # Back up to process this line normally
                                break
                            i += 1
                    else:
                        bullets.append(rest)
                        
                elif line.startswith('LEFT:'):
                    left_col = [b.strip() for b in line[5:].split('|') if b.strip()]
                    
                elif line.startswith('RIGHT:'):
                    right_col = [b.strip() for b in line[6:].split('|') if b.strip()]
                    
                elif line.startswith('SOURCES:'):
                    sources = [s.strip() for s in line[8:].split(',') if s.strip()]
                
                i += 1
            
            # Create slide if we have valid content
            if title != "Untitled Slide" and (bullets or left_col):
                slides.append(SlideContent(
                    slide_type=slide_type,
                    title=title,
                    bullet_points=bullets[:6],  # Limit to 6 bullets max
                    left_column=left_col,
                    right_column=right_col,
                    source_references=sources
                ))
        
        return slides
    
    def _create_fallback_slides(
        self, 
        request: PresentationRequest, 
        sources: List
    ) -> List[SlideContent]:
        """Create high-quality fallback slides with actual document content."""
        slides = []
        
        if not sources:
            # Minimal fallback when no sources available
            slides.append(SlideContent(
                slide_type=SlideType.BULLET_POINTS,
                title=f"{request.topic}",
                bullet_points=[
                    f"📌 Presentation topic: {request.topic}",
                    "ℹ️ No source documents found in knowledge base",
                    "💡 Upload relevant documents for better content generation",
                    "🔍 Ensure documents contain information related to this topic"
                ]
            ))
            return slides
        
        # Slide 1: Introduction with extracted key points
        intro_bullets = [f"📌 Key information about {request.topic.lower()}"]
        
        # Extract first few meaningful sentences from first source
        first_source_sentences = []
        for source in sources[:2]:
            sentences = [s.strip() + '.' for s in source.content.split('.') if len(s.strip()) > 25]
            first_source_sentences.extend(sentences[:2])
        
        intro_bullets.extend(first_source_sentences[:4])
        
        slides.append(SlideContent(
            slide_type=SlideType.BULLET_POINTS,
            title=f"{request.topic}",
            bullet_points=intro_bullets[:5]
        ))
        
        # Slides 2-N: Content from different sources
        unique_sources = list(set(s.source_file for s in sources[:5]))  # Use top 5 sources
        
        for idx, source_file in enumerate(unique_sources[:4], 2):  # Max 4 content slides
            source_chunks = [s for s in sources if s.source_file == source_file]
            
            if source_chunks:
                bullets = []
                
                # Extract meaningful paragraphs and points
                for chunk in source_chunks[:3]:  # Top 3 chunks per source
                    # Split into sentences
                    sentences = chunk.content.split('.')
                    
                    for sent in sentences:
                        sent = sent.strip()
                        # Look for substantial sentences (not too short, not too long)
                        if 30 < len(sent) < 120:
                            # Clean up the sentence
                            clean_sent = sent.strip()
                            if clean_sent and not clean_sent.startswith(('http', 'www', '...')):
                                bullets.append(clean_sent + ('.' if not clean_sent.endswith('.') else ''))
                                if len(bullets) >= 5:
                                    break
                    
                    if len(bullets) >= 5:
                        break
                
                # If we have good bullets, create a slide
                if bullets and len(bullets) >= 3:
                    # Create a descriptive title
                    title = f"Insights: {source_file.replace('_', ' ').replace('.pdf', '').replace('.txt', '')[:40]}"
                    if idx == 2:
                        title = f"Key Points from Analysis"
                    elif idx == 3:
                        title = f"Additional Findings"
                    elif idx == 4:
                        title = f"Supporting Evidence"
                    
                    slides.append(SlideContent(
                        slide_type=SlideType.BULLET_POINTS,
                        title=title,
                        bullet_points=bullets[:5],  # Max 5 bullets
                        source_references=[source_file]
                    ))
        
        # Ensure we have at least some slides
        if len(slides) < 2:
            # Add a general slide from aggregated content
            all_bullets = []
            for source in sources[:5]:
                sentences = [s.strip() + '.' for s in source.content.split('.') if 30 < len(s.strip()) < 100]
                all_bullets.extend(sentences[:2])
            
            if all_bullets:
                slides.append(SlideContent(
                    slide_type=SlideType.BULLET_POINTS,
                    title="Key Findings from Documents",
                    bullet_points=all_bullets[:5]
                ))
        
        return slides[:request.num_slides - 2 if request.num_slides > 2 else request.num_slides]  # Leave room for title and summary
    
    def _create_pptx(
        self, 
        slides: List[SlideContent], 
        theme: PresentationTheme,
        presentation_id: str
    ) -> Tuple[Path, str]:
        """Create the actual PowerPoint file."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)
        
        colors = THEME_COLORS.get(theme.value, THEME_COLORS["professional"])
        
        for slide_content in slides:
            if slide_content.slide_type == SlideType.TITLE:
                self._add_title_slide(prs, slide_content, colors)
            elif slide_content.slide_type == SlideType.TWO_COLUMN:
                self._add_two_column_slide(prs, slide_content, colors)
            elif slide_content.slide_type == SlideType.SUMMARY:
                self._add_summary_slide(prs, slide_content, colors)
            else:  # BULLET_POINTS or CONTENT
                self._add_bullet_slide(prs, slide_content, colors)
        
        # Generate filename
        safe_topic = re.sub(r'[^\w\s-]', '', slides[0].title if slides else 'presentation')[:30]
        safe_topic = safe_topic.replace(' ', '_')
        filename = f"{safe_topic}_{presentation_id}.pptx"
        filepath = self.output_dir / filename
        
        prs.save(str(filepath))
        
        return filepath, filename
    
    def _add_title_slide(self, prs, content: SlideContent, colors: Dict):
        """Add a title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add colored header bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["primary"]
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if content.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2),
                Inches(11.333), Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = content.subtitle
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = colors["secondary"]
            sub_para.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5),
            Inches(11.333), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = datetime.now().strftime("%B %d, %Y")
        date_para.font.size = Pt(16)
        date_para.font.color.rgb = colors["text"]
        date_para.alignment = PP_ALIGN.CENTER
    
    def _add_bullet_slide(self, prs, content: SlideContent, colors: Dict):
        """Add a bullet point slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Header bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        # Bullet points
        if content.bullet_points:
            body_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.7),
                Inches(11.833), Inches(5)
            )
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            
            for i, bullet in enumerate(content.bullet_points):
                if i == 0:
                    para = body_frame.paragraphs[0]
                else:
                    para = body_frame.add_paragraph()
                
                para.text = f"• {bullet}"
                para.font.size = Pt(22)
                para.font.color.rgb = colors["text"]
                para.space_after = Pt(12)
        
        # Source reference
        if content.source_references:
            src_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(7),
                Inches(12.333), Inches(0.4)
            )
            src_frame = src_box.text_frame
            src_para = src_frame.paragraphs[0]
            src_para.text = f"Source: {', '.join(content.source_references)}"
            src_para.font.size = Pt(10)
            src_para.font.italic = True
            src_para.font.color.rgb = colors["secondary"]
    
    def _add_two_column_slide(self, prs, content: SlideContent, colors: Dict):
        """Add a two-column slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Header bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        # Left column
        if content.left_column:
            left_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.7),
                Inches(5.8), Inches(5)
            )
            left_frame = left_box.text_frame
            left_frame.word_wrap = True
            
            for i, bullet in enumerate(content.left_column):
                if i == 0:
                    para = left_frame.paragraphs[0]
                else:
                    para = left_frame.add_paragraph()
                para.text = f"• {bullet}"
                para.font.size = Pt(20)
                para.font.color.rgb = colors["text"]
        
        # Right column
        if content.right_column:
            right_box = slide.shapes.add_textbox(
                Inches(6.8), Inches(1.7),
                Inches(5.8), Inches(5)
            )
            right_frame = right_box.text_frame
            right_frame.word_wrap = True
            
            for i, bullet in enumerate(content.right_column):
                if i == 0:
                    para = right_frame.paragraphs[0]
                else:
                    para = right_frame.add_paragraph()
                para.text = f"• {bullet}"
                para.font.size = Pt(20)
                para.font.color.rgb = colors["text"]
    
    def _add_summary_slide(self, prs, content: SlideContent, colors: Dict):
        """Add a summary/conclusion slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Full colored background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1),
            Inches(11.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Summary points
        if content.bullet_points:
            body_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.5),
                Inches(10.333), Inches(4)
            )
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            
            for i, bullet in enumerate(content.bullet_points):
                if i == 0:
                    para = body_frame.paragraphs[0]
                else:
                    para = body_frame.add_paragraph()
                
                para.text = f"✓ {bullet}"
                para.font.size = Pt(24)
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                para.space_after = Pt(16)
                para.alignment = PP_ALIGN.CENTER
